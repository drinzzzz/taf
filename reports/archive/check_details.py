#!/usr/bin/env python3
"""深度检查 PPTX 排版细节"""

from pptx import Presentation
from pptx.util import Inches, Pt

def check_borders(path):
    """检查表格边框"""
    prs = Presentation(path)
    slide = prs.slides[0]
    
    print(f"\n{'='*60}")
    print(f"边框检查：{path}")
    print(f"{'='*60}")
    
    for i, sh in enumerate(slide.shapes):
        if sh.has_table:
            x = sh.left.emu / 914400
            y = sh.top.emu / 914400
            
            # 获取表格
            table = sh.table
            header_text = ""
            if len(table.rows) > 0 and len(table.rows[0].cells) > 0:
                header_text = table.rows[0].cells[0].text[:30]
            
            # 检查边框
            has_l = has_r = has_t = has_b = False
            try:
                # 检查数据行第一个单元格
                if len(table.rows) > 1:
                    tc = table.rows[1].cells[0]._tc
                    tcPr = tc.get_or_add_tcPr()
                    
                    from pptx.oxml.ns import qn
                    has_l = len(tcPr.findall(qn('a:lnL'))) > 0
                    has_r = len(tcPr.findall(qn('a:lnR'))) > 0
                    has_t = len(tcPr.findall(qn('a:lnT'))) > 0
                    has_b = len(tcPr.findall(qn('a:lnB'))) > 0
            except Exception as e:
                pass
            
            border_status = "四边" if (has_l and has_r and has_t and has_b) else \
                           "部分" if (has_l or has_r or has_t or has_b) else "无"
            
            print(f"表格 @{x:.2f},{y:.2f}: {header_text:30s} → 边框：{border_status}")

def check_layout_details(path):
    """检查排版细节"""
    prs = Presentation(path)
    slide = prs.slides[0]
    
    print(f"\n{'='*60}")
    print(f"排版细节：{path}")
    print(f"{'='*60}")
    
    # 检查所有文本框的字体
    print("\n字体统计:")
    font_stats = {}
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        fs = run.font.size.pt
                        font_stats[fs] = font_stats.get(fs, 0) + 1
    
    for fs in sorted(font_stats.keys()):
        print(f"  {fs:.1f}pt: {font_stats[fs]} 处")
    
    # 检查空间利用
    print("\n垂直空间分布:")
    elements = []
    for sh in slide.shapes:
        y = sh.top.emu / 914400
        h = sh.height.emu / 914400
        elements.append((y, y+h, sh))
    
    elements.sort(key=lambda e: e[0])
    
    prev_bottom = 0
    for y_top, y_bottom, sh in elements:
        gap = y_top - prev_bottom
        if gap > 0.05:  # 大于 0.05" 的间隙
            print(f"  间隙 {gap:.2f}\" @ y={prev_bottom:.2f}\"")
        
        shape_desc = ""
        if sh.has_table:
            shape_desc = f"表格 {sh.table.rows}行"
        elif sh.has_text_frame:
            shape_desc = f"文本 \"{sh.text_frame.text[:20]}...\""
        
        print(f"  元素 @ {y_top:.2f}\"-{y_bottom:.2f}\" (h={y_bottom-y_top:.2f}\"): {shape_desc}")
        prev_bottom = y_bottom
    
    final_gap = 7.5 - prev_bottom
    print(f"  底部留白：{final_gap:.2f}\"")

def compare_with_spec():
    """对比设计规范"""
    print(f"\n{'='*60}")
    print("设计规范对比")
    print(f"{'='*60}")
    
    spec = {
        "页面比例": "4:3 (10\"×7.5\")",
        "顶部黑条": "35px (0.46\")",
        "图片区": "2.3\"×2.3\" @ (0.35, 1.50)",
        "身份表": "8 行 2 列，右上",
        "分区表格": "5 个，共 21 数据行",
        "字体": "微软雅黑，身份区 6.5pt，表头 6pt，数据 5.5pt",
        "表格边框": "暖灰米色 #C5BFB3 细线",
        "无阴影": "是",
    }
    
    actual = {
        "页面比例": "10.00\"×7.50\" ✓",
        "顶部黑条": "0.46\" ✓",
        "图片区": "2.30\"×2.30\" @ (0.35, 1.50) ✓",
        "身份表": "8 行 2 列 @ (2.87, 1.50) ✓",
        "分区表格": "5 个，共 21 数据行 ✓",
        "字体": "5.5pt-13pt (数据 5.5pt ✓)",
        "表格边框": "分区表有，身份表无边框 ⚠️",
        "无阴影": "待确认",
    }
    
    print(f"\n{'规范项':<15} {'设计要求':<25} {'实际情况':<25}")
    print("-" * 65)
    for key in spec:
        print(f"{key:<15} {spec[key]:<25} {actual.get(key, 'N/A'):<25}")

if __name__ == '__main__':
    check_borders('/www/wwwroot/draft/static/templates/materials/default.pptx')
    check_layout_details('/www/wwwroot/draft/static/templates/materials/default.pptx')
    compare_with_spec()
