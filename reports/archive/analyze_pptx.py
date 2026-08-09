#!/usr/bin/env python3
"""分析材料表 PPTX 模板排版质量"""

from pptx import Presentation
from pptx.util import Inches, Pt

SW, SH = 10.0, 7.5  # 4:3 比例

def analyze_pptx(path, label=""):
    print(f"\n{'='*60}")
    print(f"分析文件：{path} {label}")
    print(f"{'='*60}")
    
    prs = Presentation(path)
    issues = []
    
    # 检查幻灯片数量
    print(f"\n📊 幻灯片数量：{len(prs.slides)}")
    if len(prs.slides) != 1:
        issues.append(("P2", "幻灯片数量异常", f"期望 1 页，实际{len(prs.slides)}页"))
    
    slide = prs.slides[0]
    
    # 检查页面尺寸
    print(f"📐 页面尺寸：{prs.slide_width.emu/914400:.2f}\" × {prs.slide_height.emu/914400:.2f}\"")
    if abs(prs.slide_width.emu/914400 - SW) > 0.1:
        issues.append(("P1", "页面宽度不符", f"期望{SW}\", 实际{prs.slide_width.emu/914400:.2f}\""))
    if abs(prs.slide_height.emu/914400 - SH) > 0.1:
        issues.append(("P1", "页面高度不符", f"期望{SH}\", 实际{prs.slide_height.emu/914400:.2f}\""))
    
    # 收集所有形状信息
    shapes_info = []
    tables_info = []
    max_y = 0
    
    print(f"\n📋 形状分析:")
    for i, sh in enumerate(slide.shapes):
        x = sh.left.emu / 914400
        y = sh.top.emu / 914400
        w = sh.width.emu / 914400
        h = sh.height.emu / 914400
        bottom_y = y + h
        max_y = max(max_y, bottom_y)
        
        shape_type = "Unknown"
        if sh.has_text_frame:
            shape_type = "TextBox"
            text_preview = sh.text_frame.text[:50].replace('\n', ' ') if sh.text_frame.text else "(空)"
        elif sh.has_table:
            shape_type = "Table"
            tables_info.append({
                'idx': i,
                'x': x, 'y': y, 'w': w, 'h': h,
                'rows': len(sh.table.rows),
                'cols': len(sh.table.columns)
            })
            text_preview = f"表格 {len(sh.table.rows)}行×{len(sh.table.columns)}列"
        else:
            shape_type = "Shape"
            text_preview = ""
        
        shapes_info.append({
            'idx': i, 'type': shape_type,
            'x': x, 'y': y, 'w': w, 'h': h, 'bottom': bottom_y
        })
        
        print(f"  [{i:2d}] {shape_type:8s} @ ({x:5.2f}, {y:5.2f}) {w:5.2f}×{h:5.2f}  bottom={bottom_y:5.2f}  {text_preview[:40]}")
    
    # 检查内容是否超出页面
    print(f"\n⚠️  最大 Y 坐标：{max_y:.2f}\" (页面限制：{SH}\")")
    if max_y > SH:
        overflow = max_y - SH
        issues.append(("P1", "内容超出页面", f"底部超出{overflow:.2f}\""))
    elif max_y > SH - 0.1:
        issues.append(("P3", "内容接近页面边缘", f"距底部仅{SH - max_y:.2f}\""))
    
    # 分析表格
    print(f"\n📊 表格详细分析:")
    expected_sections = [
        ('■ 规格参数 SPECIFICATIONS', 5),
        ('■ 技术性能 TECHNICAL', 4),
        ('■ 商务信息 COMMERCIAL', 4),
        ('■ 施工信息 CONSTRUCTION', 6),
        ('■ 维护保养 MAINTENANCE', 2),
    ]
    
    # 预期的右上身份表：8 行 2 列
    # 预期的 5 个分区表格
    print(f"  期望：1 个身份表 (8 行×2 列) + 5 个分区表格")
    print(f"  实际找到 {len(tables_info)} 个表格")
    
    # 按 Y 坐标排序表格
    tables_sorted = sorted(tables_info, key=lambda t: t['y'])
    
    for ti, tbl in enumerate(tables_sorted):
        print(f"\n  表格 #{ti}:")
        print(f"    位置：({tbl['x']:.2f}, {tbl['y']:.2f}) 尺寸：{tbl['w']:.2f}×{tbl['h']:.2f}")
        print(f"    行列：{tbl['rows']}行 × {tbl['cols']}列")
        
        # 获取实际表格对象
        for sh in slide.shapes:
            if sh.has_table:
                x = sh.left.emu / 914400
                y = sh.top.emu / 914400
                if abs(x - tbl['x']) < 0.1 and abs(y - tbl['y']) < 0.1:
                    table = sh.table
                    # 检查第一行（标题行）内容
                    if len(table.rows) > 0 and len(table.rows[0].cells) > 0:
                        header_text = table.rows[0].cells[0].text[:60]
                        print(f"    标题：{header_text}")
                    
                    # 检查字体大小
                    font_sizes = set()
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if run.font.size:
                                            font_sizes.add(run.font.size.pt)
                    print(f"    字体大小：{sorted(font_sizes)} pt")
                    
                    # 检查边框
                    has_border = False
                    try:
                        tc = table.rows[1].cells[0]._tc if len(table.rows) > 1 else table.rows[0].cells[0]._tc
                        tcPr = tc.get_or_add_tcPr()
                        if tcPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lnL'):
                            has_border = True
                    except:
                        pass
                    print(f"    边框：{'有' if has_border else '无/未知'}")
                    break
    
    # 检查身份表空间利用
    print(f"\n🔍 身份表空间利用率分析:")
    # 身份表应该在右上角，与图片框同高
    # IMG_Y=1.50, IMG_SZ=2.3, 所以身份表应该在 y=1.50, 高度约 2.3
    id_table = None
    for tbl in tables_sorted:
        # 身份表特征：在右上区域，8 行
        if tbl['y'] < 2.0 and tbl['rows'] == 8:
            id_table = tbl
            break
    
    if id_table:
        print(f"  身份表高度：{id_table['h']:.2f}\"")
        print(f"  身份表行数：{id_table['rows']}行")
        # 预期 8 行充分利用 2.3" 高度
        avg_row_h = id_table['h'] / id_table['rows']
        print(f"  平均每行高度：{avg_row_h:.3f}\"")
        if avg_row_h > 0.35:
            issues.append(("P3", "身份表行高过大", f"平均每行{avg_row_h:.3f}\"，可能有留白"))
    else:
        issues.append(("P2", "未找到身份表", "期望 8 行 2 列的身份信息表"))
    
    # 检查分区表格完整性
    print(f"\n🔍 分区表格完整性检查:")
    section_tables = [t for t in tables_sorted if t['rows'] >= 3]  # 至少标题行 +2 数据行
    print(f"  找到 {len(section_tables)} 个分区表格")
    
    found_sections = []
    for tbl in section_tables:
        for sh in slide.shapes:
            if sh.has_table:
                x = sh.left.emu / 914400
                y = sh.top.emu / 914400
                if abs(x - tbl['x']) < 0.1 and abs(y - tbl['y']) < 0.1:
                    table = sh.table
                    if len(table.rows) > 0:
                        header = table.rows[0].cells[0].text if len(table.rows[0].cells) > 0 else ""
                        found_sections.append((header, len(table.rows) - 1))  # -1 for header
                        print(f"    - {header[:40]}: {len(table.rows)-1} 数据行")
                    break
    
    # 对比预期
    for exp_title, exp_rows in expected_sections:
        found = False
        for found_title, found_rows in found_sections:
            if exp_title.split()[0] in found_title or exp_title.split()[1][:4] in found_title:
                found = True
                if found_rows != exp_rows:
                    issues.append(("P2", f"字段数量不符", f"{exp_title}: 期望{exp_rows}行，实际{found_rows}行"))
                break
        if not found:
            issues.append(("P2", "分区缺失", f"未找到：{exp_title}"))
    
    # 检查黑条
    print(f"\n🔍 顶部黑条检查:")
    bar_found = False
    for sh in shapes_info:
        if sh['y'] < 0.1 and sh['h'] > 0.3:  # 顶部，高度约 0.46
            bar_found = True
            print(f"  黑条：({sh['x']:.2f}, {sh['y']:.2f}) {sh['w']:.2f}×{sh['h']:.2f}")
            # 检查文字
            for s in slide.shapes:
                x = s.left.emu / 914400
                y = s.top.emu / 914400
                if abs(x - sh['x']) < 0.1 and abs(y - sh['y']) < 0.1 and s.has_text_frame:
                    text = s.text_frame.text[:80]
                    print(f"  黑条文字：{text}")
                    break
            break
    if not bar_found:
        issues.append(("P2", "顶部黑条缺失", "期望 35px 黑色横条"))
    
    # 检查图片占位框
    print(f"\n🔍 图片占位框检查:")
    img_found = False
    for sh in shapes_info:
        # 图片框应该在左上，约 2.3×2.3
        if 0.2 < sh['x'] < 1.0 and 1.0 < sh['y'] < 2.0 and 2.0 < sh['w'] < 2.6:
            img_found = True
            print(f"  图片框：({sh['x']:.2f}, {sh['y']:.2f}) {sh['w']:.2f}×{sh['h']:.2f}")
            break
    if not img_found:
        issues.append(("P2", "图片占位框缺失或位置异常", "期望左上 2.3\"×2.3\" 灰底占位框"))
    
    # 检查底部备注区
    print(f"\n🔍 底部备注区检查:")
    bottom_elements = [sh for sh in shapes_info if sh['bottom'] > SH - 0.5]
    print(f"  底部区域元素：{len(bottom_elements)} 个")
    for sh in bottom_elements:
        print(f"    @ y={sh['y']:.2f}, bottom={sh['bottom']:.2f}")
    
    if len(bottom_elements) < 2:
        issues.append(("P2", "底部备注区不完整", "期望设计师备注 + 业主确认两个元素"))
    
    # 检查字体大小可读性
    print(f"\n🔍 字体大小检查:")
    all_font_sizes = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        all_font_sizes.append(run.font.size.pt)
    
    if all_font_sizes:
        min_fs = min(all_font_sizes)
        max_fs = max(all_font_sizes)
        avg_fs = sum(all_font_sizes) / len(all_font_sizes)
        print(f"  字体大小范围：{min_fs:.1f}pt - {max_fs:.1f}pt (平均：{avg_fs:.1f}pt)")
        
        if min_fs < 5:
            issues.append(("P3", "字体过小", f"最小字体{min_fs:.1f}pt，可能影响可读性"))
        elif min_fs < 5.5:
            issues.append(("P3", "字体偏小", f"最小字体{min_fs:.1f}pt，接近可读性下限"))
    
    # 总结
    print(f"\n{'='*60}")
    print(f"问题汇总：{len(issues)} 个")
    for severity, title, desc in sorted(issues, key=lambda x: x[0]):
        print(f"  [{severity}] {title}: {desc}")
    
    return issues

if __name__ == '__main__':
    issues_template = analyze_pptx('/www/wwwroot/draft/static/templates/materials/default.pptx', '(模板)')
    issues_sample = analyze_pptx('/www/wwwroot/draft/static/templates/materials/default_sample.pptx', '(样本)')
    
    print(f"\n{'='*60}")
    print("最终问题清单（合并去重）")
    print(f"{'='*60}")
    
    all_issues = set()
    for issue in issues_template + issues_sample:
        all_issues.add((issue[0], issue[1], issue[2]))
    
    for severity, title, desc in sorted(all_issues, key=lambda x: x[0]):
        print(f"[{severity}] {title}")
        print(f"      → {desc}")
