#!/usr/bin/env python3
"""TAF 点位符号总览 PPT — 两页:
页1 白色描边勾勒; 页2 白色实心 + 内部细节镂空
图标 path 整条解析(支持混合 M/m), 按跳跃断点拆分折线组, 避免假连线与负坐标
"""
import sys, json, html as _html
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

MAKI = json.load(open('/root/TAF/frontend/maki_symbols.json', encoding='utf-8'))['symbols']
items = sorted(MAKI.keys())

EMU = 9525.0
PAGE_W, PAGE_H = 3840, 2160
prs = Presentation()
prs.slide_width = Emu(int(PAGE_W * EMU))
prs.slide_height = Emu(int(PAGE_H * EMU))

def no_shadow(obj):
    try:
        el = obj._element
        st = el.find(qn('p:style'))
        if st is not None:
            el.remove(st)
        for e2 in el.spPr.findall(qn('a:effectLst')):
            el.spPr.remove(e2)
    except Exception:
        pass

def set_font(r, size_pt, color=RGBColor(0x11, 0x11, 0x11), bold=False):
    r.font.size = Pt(size_pt)
    r.font.color.rgb = color
    r.font.name = 'MiSans'
    r.font.bold = bold
    try:
        rPr = r._r.get_or_add_rPr()
        for tag in ('a:ea', 'a:cs'):
            el = rPr.makeelement(qn(tag), {'typeface': 'MiSans'})
            rPr.append(el)
    except Exception:
        pass

def add_text(slide, cx, cy, txt, size_pt, bold=False, w_px=900, color=RGBColor(0x11, 0x11, 0x11)):
    tb = slide.shapes.add_textbox(Emu(int((cx - w_px / 2) * EMU)), Emu(int(cy * EMU)),
                                  Emu(int(w_px * EMU)), Emu(int((size_pt * 1.6 + 8) * EMU)))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 2
    r = p.add_run()
    r.text = txt
    set_font(r, size_pt, bold=bold, color=color)
    no_shadow(tb)

def path_to_lines(d):
    """整条解析 (svgpathtools 正确处理混合 M/m 相对命令), 逐段采样,
    相邻采样点跳跃(>1.5 unit)即断开 → 折线组, 附带隐闭判定(首尾近=闭合)"""
    from svgpathtools import parse_path
    p = parse_path(_html.unescape(d))
    lines = []
    cur = []
    prev = None
    for seg in p:
        n = max(2, int(seg.length() / 0.9) + 1)
        for i in range(n + 1):
            z = seg.point(i / n)
            pt = (z.real, z.imag)
            if prev is not None and cur:
                d0 = ((pt[0] - prev[0]) ** 2 + (pt[1] - prev[1]) ** 2) ** 0.5
                if d0 > 1.5:               # 跳跃 → 断线
                    if len(cur) >= 2:
                        lines.append(cur)
                    cur = []
            cur.append(pt)
            prev = pt
    if len(cur) >= 2:
        lines.append(cur)
    out = []
    for ln in lines:
        x0, y0 = ln[0]; x1, y1 = ln[-1]
        out.append({'pts': ln, 'close': ((x0 - x1) ** 2 + (y0 - y1) ** 2) ** 0.5 < 0.8})
    return out

def bbox(pts):
    xs = [p[0] for p in pts]; ys = [p[1] for p in pts]
    return min(xs), min(ys), max(xs), max(ys)

def add_freeform(slide, pts_px, close, fill_rgb=None, line_rgb=None, lw_px=0):
    p0 = pts_px[0]
    fb = slide.shapes.build_freeform(int(p0[0]), int(p0[1]), scale=1.0)
    fb.add_line_segments([(int(a), int(b)) for a, b in pts_px[1:]], close=close)
    sp = fb.convert_to_shape()
    if fill_rgb is not None:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill_rgb
    else:
        sp.fill.background()
    if line_rgb is not None:
        sp.line.color.rgb = line_rgb
        sp.line.width = Emu(int(9525 * lw_px))
    else:
        sp.line.fill.background()
    no_shadow(sp)
    return sp

def draw_page(mode):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if mode == 'outline':
        TITLE = '兴顺里 TAF · 点位符号总览（18 项）· 白色描边'
        SUB = '黑色正圆 = 设施图标 · 白色轮廓线 = 布点符号'
    else:
        TITLE = '兴顺里 TAF · 点位符号总览（18 项）· 实心 + 镂空细节'
        SUB = '黑色正圆 = 设施图标 · 白色实心 + 内部细节镂空'
    GRID_TOP = 210
    COLS, ROWS = 4, 5
    CELL_W, CELL_H = PAGE_W / COLS, (PAGE_H - GRID_TOP) / ROWS
    CIRCLE_D = 150
    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    add_text(slide, PAGE_W / 2, 60, TITLE, 34, bold=True, w_px=2800)
    add_text(slide, PAGE_W / 2, 130, SUB, 18, w_px=2000)
    for idx, item in enumerate(items):
        sym = MAKI[item]
        col, row = idx % COLS, idx // COLS
        cx = col * CELL_W + CELL_W / 2
        cy_c = GRID_TOP + row * CELL_H + CELL_H * 0.30
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Emu(int((cx - CIRCLE_D / 2) * EMU)),
                                        Emu(int((cy_c - CIRCLE_D / 2) * EMU)),
                                        Emu(int(CIRCLE_D * EMU)), Emu(int(CIRCLE_D * EMU)))
        circle.fill.solid(); circle.fill.fore_color.rgb = BLACK
        circle.line.fill.background()
        no_shadow(circle)
        circle.name = f'circ_{item}'
        k = (CIRCLE_D * 0.86) / 15.0
        lines = path_to_lines(sym['path'])
        if not lines:
            continue
        tx0 = min(min(bbox(l['pts'])[0] for l in lines), 0)
        ty0 = min(min(bbox(l['pts'])[1] for l in lines), 0)
        tx1 = max(max(bbox(l['pts'])[2] for l in lines), 15)
        ty1 = max(max(bbox(l['pts'])[3] for l in lines), 15)
        tarea = max(1e-6, (tx1 - tx0) * (ty1 - ty0))
        def to_px(u, v):
            return ((cx + (u - (tx0 + tx1) / 2) * k) * EMU, (cy_c + (v - (ty0 + ty1) / 2) * k) * EMU)
        for l in lines:
            b = bbox(l['pts'])
            cxm = (b[0] + b[2]) / 2; cym = (b[1] + b[3]) / 2
            l['inner'] = (l['close'] and (b[2] - b[0]) * (b[3] - b[1]) < 0.30 * tarea
                          and tx0 < cxm < tx1 and ty0 < cym < ty1)
        if mode == 'outline':
            for l in lines:
                px = [to_px(u, v) for u, v in l['pts']]
                add_freeform(slide, px, l['close'], fill_rgb=None, line_rgb=WHITE, lw_px=2.0)
        else:
            for l in lines:
                px = [to_px(u, v) for u, v in l['pts']]
                add_freeform(slide, px, l['close'], fill_rgb=WHITE, line_rgb=None)
            for l in lines:
                if not l['inner']:
                    continue
                px = [to_px(u, v) for u, v in l['pts']]
                if l['close']:
                    add_freeform(slide, px, True, fill_rgb=BLACK, line_rgb=None)
                else:
                    add_freeform(slide, px, False, fill_rgb=None, line_rgb=BLACK, lw_px=2.5)
        t0 = cy_c + CIRCLE_D / 2 + 26
        add_text(slide, cx, t0, item, 26, bold=True)
        add_text(slide, cx, t0 + 52, sym['name_zh'], 20)
        add_text(slide, cx, t0 + 100, sym['layer'], 12, w_px=880)

draw_page('outline')
draw_page('solid')
out = '/data/disk1/wwwroot/taf/exports_dev/XING_SHUN_LI_TAF_Symbols.pptx'
prs.save(out)
print('saved', out, '| pages 2 | items', len(items))
