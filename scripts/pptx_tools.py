#!/usr/bin/env python3
"""pptx 高层工具(封装多轮踩坑的正确做法) — 供 TAF 成果 PPT 等日常修改使用
要点: 一切段落属性走 python-pptx API 或 a:pPr(禁止 p:pPr);
      graphicFrame/组的 xfrm 定位用正则; 页面归属按页眉 60pt 标题的标准项编号."""
import re
from pptx import Presentation
from pptx.util import Pt, Emu
from lxml import etree

P_NS = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
A_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
EMU = 9525.0
FONT = 'MiSans'
SEC_PT, BODY_PT = 24.0, 20.0  # 说明区默认规格(第5页模板)

# ---------- 基础 ----------
def load(path):
    return Presentation(path)

def save(prs, path):
    prs.save(path)

def x0_of(el):
    m = re.search(r'<(?:a|p):xfrm[^>]*>\s*<(?:a|p):off x="(-?\d+)"', el.xml)
    return int(m.group(1)) if m else None

def top_of(el):
    m = re.search(r'<(?:a|p):xfrm[^>]*>\s*<(?:a|p):off x="-?\d+" y="(-?\d+)"', el.xml)
    return int(m.group(1)) if m else None

def page_item(slide):
    """页眉 60pt 大标题中的标准项编号(与页码无关)"""
    for sh in slide.shapes:
        if not sh.has_text_frame or sh.top is None or sh.top / EMU > 160:
            continue
        for pa in sh.text_frame.paragraphs:
            for r in pa.runs:
                if r.font.size and abs(r.font.size.pt - 60) < 1:
                    m = re.search(r'\b(P\d-\d{2})\b', sh.text_frame.text)
                    if m:
                        return m.group(1)
    return None

def clear_region(slide, x_min_px):
    """删除页面右/指定区(x≥x_min_px)的顶层 sp/cxnSp/graphicFrame(表), 防叠加"""
    spTree = slide.shapes._spTree
    n = 0
    for child in list(spTree):
        if child.tag in (P_NS + 'sp', P_NS + 'cxnSp', P_NS + 'graphicFrame'):
            x0 = x0_of(child)
            if x0 is not None and x0 >= x_min_px * EMU:
                spTree.remove(child)
                n += 1
    return n

def no_shadow(el):
    """去阴影(图形新增后必调): 删 p:style + spPr effectLst"""
    st = el.find(P_NS + 'style')
    if st is not None:
        el.remove(st)
    spPr = el.find(P_NS + 'spPr')
    if spPr is not None:
        for e in spPr.findall(A_NS + 'effectLst'):
            spPr.remove(e)

# ---------- 段落/文本 ----------
def set_run(r, text, size_pt, bold=False, rgb=(0x26, 0x29, 0x2E), font=FONT):
    r.text = text
    r.font.size = Pt(size_pt)
    r.font.bold = bold
    r.font.color.rgb = _rgb(rgb)
    if font:
        r.font.name = font
        rPr = r._r.get_or_add_rPr()
        for tag in ('latin', 'ea', 'cs'):
            e = rPr.find(A_NS + tag)
            if e is None:
                e = etree.SubElement(rPr, A_NS + tag)
            e.set('typeface', font)

def _rgb(t):
    from pptx.dml.color import RGBColor
    return RGBColor(*t)

def para_fmt(pa, line=None, before_pt=None, after_pt=None):
    """段落行距/段前/段后 — 一律用 python API(写 a:pPr 合法)"""
    if line is not None:
        pa.line_spacing = line
    if before_pt is not None:
        pa.space_before = Pt(before_pt)
    if after_pt is not None:
        pa.space_after = Pt(after_pt)

def para_bullet(pa, char='▪'):
    """内置项目符号(方块) — 注入 a:pPr 尾部(顺序: buFont 后 buChar)"""
    pPr = pa._p.find(A_NS + 'pPr')
    if pPr is None:
        pPr = pa._p.makeelement(A_NS + 'pPr', {})
        pa._p.insert(0, pPr)
    pPr.set('marL', '370840')
    pPr.set('indent', '-370840')
    if pPr.find(A_NS + 'buFont') is None:
        etree.SubElement(pPr, A_NS + 'buFont').set('typeface', 'Arial')
    if pPr.find(A_NS + 'buChar') is None:
        etree.SubElement(pPr, A_NS + 'buChar').set('char', char)

def para_clean(pa):
    """移除错误 p:pPr(若历史文件被旧脚本污染)"""
    bad = pa._p.find(P_NS + 'pPr')
    if bad is not None:
        pa._p.remove(bad)
        return True
    return False

def first_run(pa):
    for r in pa.runs:
        if r.text.strip():
            return r
    return None

def is_sec_para(pa, min_pt=24):
    r = first_run(pa)
    return bool(r and r.font.bold and (r.font.size is None or r.font.size.pt >= min_pt))

def para_text(pa):
    return ''.join(r.text for r in pa.runs).strip()

# ---------- 文本行高/折行估算(说明区排版) ----------
def half_units(t):
    return sum(2 if ord(ch) > 0x2E7F else 1 for ch in t)

def lines_of(text, pt, width_px):
    cap = width_px / (pt * 1.35)
    return max(1, math.ceil(half_units(text) / 2 / cap))

def textbox(slide, x, y, w, lines_est, h_min=40):
    """建文本框(可拉伸: word_wrap on, auto_size off, 高按行估算)"""
    from pptx.enum.text import MSO_AUTO_SIZE
    tb = slide.shapes.add_textbox(Emu(int(x * EMU)), Emu(int(y * EMU)),
                                  Emu(int(w * EMU)), Emu(int(max(h_min, lines_est * 34 + 8) * EMU)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    no_shadow(tb._element)
    return tb, tf

def add_para(tf, text, size_pt, bold=False, rgb=(0x26, 0x29, 0x2E), font=FONT,
             line=None, before=None, after=None, bullet=None, first=False):
    pa = tf.paragraphs[0] if first else tf.add_paragraph()
    para_fmt(pa, line, before, after)
    if bullet:
        para_bullet(pa, bullet)
    r = pa.add_run()
    set_run(r, text, size_pt, bold, rgb, font)
    return pa

# ---------- 组变换坐标 ----------
def grp_geo(gel):
    xf = gel.find(P_NS + 'grpSpPr/' + A_NS + 'xfrm')
    if xf is None:
        return None
    Go = xf.find(A_NS + 'off'); Ge = xf.find(A_NS + 'ext')
    Gc = xf.find(A_NS + 'chOff'); Gce = xf.find(A_NS + 'chExt')
    return (int(Go.get('x')), int(Go.get('y')),
            int(Gc.get('x')), int(Gc.get('y')),
            int(Ge.get('cx')) / int(Gce.get('cx')),
            int(Ge.get('cy')) / int(Gce.get('cy')))

def to_display(x, y, geo):
    ogx, ogy, cgx, cgy, sx, sy = geo
    return ogx + (x - cgx) * sx, ogy + (y - cgy) * sy

import math
