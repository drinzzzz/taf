#!/usr/bin/env python3
"""TAF 完整版 PPTX — 27 页分页 (与 PNG/SVG/EMF 图系同清单)
页构成与同名图一致:
  composite / basemap          : basemap-clean.png 垫底(灰线细节) + 语义矢量 (+composite 含全部点位矢量)
  basemap-clean                : basemap-clean.png 垫底 + 图框矢量
  map-TAF-{6 层}               : 该语义层 + 图框矢量 (与 png 同内容, 无灰线)
  layer-Px-0X (18)             : 图框 + 该标准项点位符号+编号 (与 png 同内容)
映射: 页 3840×2160, 内边距 2px (与 PNG/SVG 精确一致)
"""
import sys, os, math, json
sys.path.insert(0, '/root/TAF/scripts')
import export_figures as ef
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

ents = ef.load_dxf()
ef.maki = json.load(open(ef.MAKI_JSON, encoding='utf-8'))['symbols']
ef.placements = json.load(open(ef.PLACE_JSON, encoding='utf-8'))
placements = ef.placements

SEM = {'TAF-DRAWING_BORDER', 'TAF-BOUNDARY', 'TAF-BUILDING', 'TAF-CHANNEL',
       'TAF-NODE', 'TAF-GREEN', 'TAF-ROAD', 'TAF-FACADE'}
geo_all = [e for e in ents if e[0] in SEM]

bb = ef.border_bbox(geo_all)
minx, miny, maxx, maxy = bb
EMU = 9525.0
PAGE_W, PAGE_H = 3840, 2160
PAD = 2.0
scale = (PAGE_H - 2 * PAD) / (maxy - miny)
cw = (maxx - minx) * scale
ox = (PAGE_W - cw) / 2
oy = PAD

def X(x): return (ox + (x - minx) * scale) * EMU
def Y(y): return (oy + (maxy - y) * scale) * EMU

def hex2rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def no_shadow(obj):
    try:
        el = obj._element
        st = el.find(qn('p:style'))
        if st is not None:
            el.remove(st)
        for e2 in el.spPr.findall(qn('a:effectLst')):
            el.spPr.remove(e2)
    except Exception:
        pass

BG = '/data/disk1/wwwroot/taf/exports_dev/basemap-clean.png'
BG_SIZE = None
if os.path.exists(BG):
    from PIL import Image as _I
    BG_SIZE = _I.open(BG).size

CAT_COLOR = {'P1': '#4f8cff', 'P2': '#67c23a', 'P3': '#e6a23c', 'P4': '#c97dff', 'P5': '#f56c6c', 'P6': '#8b8fa3'}

def add_bg(slide):
    if not BG_SIZE:
        return
    iw, ih = BG_SIZE
    pic = slide.shapes.add_picture(BG, Emu(int((ox - PAD) * 9525)), Emu(int(0)),
                                   Emu(int(iw * 9525)), Emu(int(ih * 9525)))
    pic.name = 'BG_clean'
    spTree = slide.shapes._spTree
    el = pic._element
    spTree.remove(el)
    spTree.insert(2, el)
    no_shadow(pic)

def add_poly_geo(slide, layer, typ, pl, lw_extra=1.0):
    col = ef.LAYER_COLORS.get(layer, ef.LAYER_COLORS.get(layer.split('-')[-1], '#9aa3b2'))
    lwpx = max(1.0, ef.LAYER_LW.get(layer, 1.6))
    if typ == 'poly':
        pts, closed = pl
        if len(pts) < 2:
            return
        sp_pts = [(X(a), Y(b)) for a, b in pts]
        fb = slide.shapes.build_freeform(int(sp_pts[0][0]), int(sp_pts[0][1]), scale=1.0)
        fb.add_line_segments([(int(a), int(b)) for a, b in sp_pts[1:]], close=bool(closed))
        sp = fb.convert_to_shape()
        sp.fill.background()
        sp.line.color.rgb = hex2rgb(col)
        sp.line.width = Emu(int(9525 * lwpx * lw_extra))
        no_shadow(sp)
        sp.name = layer + '_' + ('poly' if closed else 'line')
    elif typ == 'line':
        (x1, y1, x2, y2) = pl
        sp_pts = [(X(x1), Y(y1)), (X(x2), Y(y2))]
        fb = slide.shapes.build_freeform(int(sp_pts[0][0]), int(sp_pts[0][1]), scale=1.0)
        fb.add_line_segments([(int(sp_pts[1][0]), int(sp_pts[1][1]))], close=False)
        sp = fb.convert_to_shape()
        sp.fill.background()
        sp.line.color.rgb = hex2rgb(col)
        sp.line.width = Emu(int(9525 * lwpx * lw_extra))
        no_shadow(sp)
        sp.name = layer + '_line'

def add_maki_pt(slide, x, y, path_d, color, size_px, seq=None, base=''):
    pts_w = _path_polylines(path_d)   # 相对 15box, 中心 ~7.5
    if len(pts_w) < 3:
        return
    k = size_px / 15.0
    cx, cy = X(x), Y(y)
    px_pts = [(cx + (u - 7.5) * k * EMU, cy + (v - 7.5) * k * EMU) for u, v in pts_w]
    fb = slide.shapes.build_freeform(int(px_pts[0][0]), int(px_pts[0][1]), scale=1.0)
    fb.add_line_segments([(int(a), int(b)) for a, b in px_pts[1:]], close=True)
    sp = fb.convert_to_shape()
    sp.fill.background()
    sp.line.color.rgb = hex2rgb(color)
    sp.line.width = Emu(int(9525 * 1.2))
    no_shadow(sp)
    sp.name = f'PT_{base}'
    # 编号
    if base:
        txt = base if seq is None or seq <= 1 else f'{base}-{seq}'
        tb = slide.shapes.add_textbox(Emu(int(cx + size_px * 0.6)), Emu(int(cy + size_px * 0.5)),
                                      Emu(int(140 * 9525)), Emu(int(24 * 9525)))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = txt
        r.font.size = Pt(9.5)
        r.font.color.rgb = hex2rgb(color)
        no_shadow(tb)
        tb.name = 'PTL_' + base

def _path_polylines(path_d):
    import html as _html
    try:
        from svgpathtools import parse_path
        p = parse_path(_html.unescape(path_d))   # 🔴 实体 &#xA;/&#x9; → 换行/制表, 否则 parse ValueError
        pts = []
        for seg in p:
            n = max(2, int(seg.length() / 0.9) + 1)
            for i in range(n):
                z = seg.point(i / n)
                pts.append((z.real, z.imag))
        return pts
    except Exception:
        return []

def add_points(slide, only_item=None):
    for fac in placements['facilities']:
        item = fac['standard_item_id']
        if only_item and item != only_item:
            continue
        sym = ef.maki.get(item)
        color = CAT_COLOR.get(item.split('-')[0], '#ffffff')
        base = item.replace('P', '').replace('-', '')
        for pl in fac['placements']:
            x, y = float(pl['x']), float(pl['y'])
            if sym and sym.get('path'):
                add_maki_pt(slide, x, y, sym['path'], color, 30, pl['seq'], base)
            else:
                # 圆 fallback
                from pptx.enum.shapes import MSO_SHAPE
                rpx = 12
                cx, cy = X(x), Y(y)
                sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - rpx * 9525)), Emu(int(cy - rpx * 9525)),
                                            Emu(int(2 * rpx * 9525)), Emu(int(2 * rpx * 9525)))
                sh.fill.background()
                sh.line.color.rgb = hex2rgb(color)
                no_shadow(sh)

def add_title(slide, idx, name):
    tb = slide.shapes.add_textbox(Emu(int(20 * 9525)), Emu(int(12 * 9525)),
                                  Emu(int(900 * 9525)), Emu(int(26 * 9525)))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    r = tf.paragraphs[0].add_run()
    r.text = f'{idx}/27  {name}'
    r.font.size = Pt(10)
    r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    no_shadow(tb)

# ── 页面清单 ──
pages = [('composite', '全部底图 + 全部点位', 'comp'),
         ('basemap', '全部底图 (无点位)', 'base'),
         ('basemap-clean', '素色底图 (无彩色)', 'clean')]
for L in ['TAF-BOUNDARY', 'TAF-BUILDING', 'TAF-CHANNEL', 'TAF-NODE', 'TAF-GREEN', 'TAF-FACADE']:
    pages.append((f'map-{L}', f'语义层 {L}', 'map:' + L))
items = [f['standard_item_id'] for f in placements['facilities'] if f['placements']]
for it in items:
    pages.append((f'layer-{it}', f'点位层 {it}', 'layer:' + it))

prs = Presentation()
prs.slide_width = Emu(int(PAGE_W * EMU))
prs.slide_height = Emu(int(PAGE_H * EMU))

for idx, (tag, title, kind) in enumerate(pages, start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if kind in ('comp', 'base', 'clean'):
        add_bg(slide)                      # 灰线细节垫底
    if kind.startswith('map:'):
        layer = kind.split(':')[1]
        for layer2, typ, pl in geo_all:
            if layer2 in (layer, 'TAF-DRAWING_BORDER'):
                add_poly_geo(slide, layer2, typ, pl)
    elif kind in ('comp', 'base'):
        for layer2, typ, pl in geo_all:
            if layer2 == 'TAF-DRAWING_BORDER' or layer2 != 'TAF-DRAWING_BORDER':
                add_poly_geo(slide, layer2, typ, pl)
    elif kind == 'clean':
        # basemap-clean 垫底已含全部内容 (灰线+框+编号), 补图框矢量供微调
        for layer2, typ, pl in geo_all:
            if layer2 == 'TAF-DRAWING_BORDER':
                add_poly_geo(slide, layer2, typ, pl)
    if kind.startswith('layer:'):
        item = kind.split(':')[1]
        for layer2, typ, pl in geo_all:
            if layer2 == 'TAF-DRAWING_BORDER':
                add_poly_geo(slide, layer2, typ, pl)
        add_points(slide, only_item=item)
    if kind == 'comp':
        add_points(slide)
    add_title(slide, idx, title)

out = '/data/disk1/wwwroot/taf/exports_dev/XING_SHUN_LI_TAF_27pages.pptx'
prs.save(out)
print('saved', out, 'pages:', len(prs.slides.__iter__.__self__._sldIdLst))
