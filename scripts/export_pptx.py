#!/usr/bin/env python3
"""TAF PPTX 生产导出 — basemap.pptx (v2: 经 LO 兼容性验证的样式组合)
页面 3840×2160; 语义几何+图框+建筑编号; 无阴影; 原生可编辑形状
"""
import sys, json
sys.path.insert(0, '/root/TAF/scripts')
import export_figures as ef
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

ents = ef.load_dxf()
SEM = {'TAF-DRAWING_BORDER', 'TAF-BOUNDARY', 'TAF-BUILDING', 'TAF-CHANNEL',
       'TAF-NODE', 'TAF-GREEN', 'TAF-ROAD', 'TAF-FACADE'}
geo = [e for e in ents if e[0] in SEM]
mtxt = [e for e in ents if e[0] == 'TAF-BUILDING_NUMBER' and e[1] == 'mtext']
bb = ef.border_bbox(geo)
minx, miny, maxx, maxy = bb
EMU = 9525.0
PAGE_W, PAGE_H = 3840, 2160
PAD = 2.0          # 🔴 与 SVG/PNG 相同的 2px 内边距 (图框外缘距画布边, 消除尺寸差)
scale = (PAGE_H - 2 * PAD) / (maxy - miny)
cw = (maxx - minx) * scale
ox = (PAGE_W - cw) / 2
oy = PAD

def X(x): return (ox + (x - minx) * scale) * EMU
def Y(y): return (oy + (maxy - y) * scale) * EMU

def hex2rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def no_shadow(obj):
    # 🔴 阴影根因: 形状带 <p:style><a:effectRef idx=N> 引用主题效果样式(默认含 outerShdw)
    #   → 即使 spPr 无 effectLst, PPT 仍按主题渲染阴影。修复: 删除 p:style 引用 + 清空 effectLst
    try:
        el = obj._element
        st = el.find(qn('p:style'))
        if st is not None:
            el.remove(st)
        for e2 in el.spPr.findall(qn('a:effectLst')):
            el.spPr.remove(e2)
    except Exception:
        pass

prs = Presentation()
prs.slide_width = Emu(int(PAGE_W * EMU))
prs.slide_height = Emu(int(PAGE_H * EMU))
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 🔴 内置素色底图 (basemap-clean.png 2160 高基准) 置于最底层 — 与矢量几何同映射, 无需手动对齐
import os as _os
BG = '/data/disk1/wwwroot/taf/exports_dev/basemap-clean.png'
if _os.path.exists(BG):
    from PIL import Image as _I
    iw, ih = _I.open(BG).size           # ≈2824×2160
    pic = slide.shapes.add_picture(BG, Emu(int((ox - PAD) * 9525)), Emu(int(0)), Emu(int(iw * 9525)), Emu(int(ih * 9525)))
    pic.name = 'BG_basemap_clean'
    # 移到 spTree 最底层 (紧跟 grpSpPr 之后)
    spTree = slide.shapes._spTree
    el = pic._element
    spTree.remove(el)
    spTree.insert(2, el)
    try:
        for el2 in pic._element.spPr.findall(qn('a:effectLst')):
            pic._element.spPr.remove(el2)
    except Exception:
        pass

cnt = {'poly': 0, 'line': 0}
for layer, typ, pl in geo:
    col = ef.LAYER_COLORS.get(layer, ef.LAYER_COLORS.get(layer.split('-')[-1], '#9aa3b2'))
    lwpx = max(1.0, (ef.LAYER_LW.get(layer, 1.6)))
    if typ == 'poly':
        pts, closed = pl
        if len(pts) < 2:
            continue
        sp_pts = [(X(a), Y(b)) for a, b in pts]
        fb = slide.shapes.build_freeform(int(sp_pts[0][0]), int(sp_pts[0][1]), scale=1.0)
        fb.add_line_segments([(int(a), int(b)) for a, b in sp_pts[1:]], close=bool(closed))
        sp = fb.convert_to_shape()
        sp.fill.background()
        sp.line.color.rgb = hex2rgb(col)
        sp.line.width = Emu(int(9525 * lwpx))
        no_shadow(sp)
        sp.name = f'{layer}_' + ('poly' if closed else 'line')
        cnt['poly' if closed else 'line'] += 1
    elif typ == 'line':
        (x1, y1, x2, y2) = pl
        sp_pts = [(X(x1), Y(y1)), (X(x2), Y(y2))]
        fb = slide.shapes.build_freeform(int(sp_pts[0][0]), int(sp_pts[0][1]), scale=1.0)
        fb.add_line_segments([(int(sp_pts[1][0]), int(sp_pts[1][1]))], close=False)
        sp = fb.convert_to_shape()
        sp.fill.background()
        sp.line.color.rgb = hex2rgb(col)
        sp.line.width = Emu(int(9525 * lwpx))
        no_shadow(sp)
        sp.name = f'{layer}_line'
        cnt['line'] += 1

for layer, typ, pl in mtxt:
    txt = ef.clean_mtext(pl[2])
    if not txt:
        continue
    tb = slide.shapes.add_textbox(Emu(int(X(pl[0]))), Emu(int(Y(pl[1]))),
                                  Emu(int(2000)), Emu(int(300)))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = txt
    r.font.size = Emu(int(10 * 9525))
    r.font.color.rgb = RGBColor(0x8A, 0x93, 0xA5)
    no_shadow(tb)
    tb.name = 'BN_text'

out = '/data/disk1/wwwroot/taf/exports_dev/basemap.pptx'
prs.save(out)
print('saved', out, '| poly', cnt['poly'], 'line', cnt['line'], 'text', len(mtxt))
